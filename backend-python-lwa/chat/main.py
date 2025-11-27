"""
Lambda handler FastAPI avec Lambda Web Adapter pour streaming Bedrock
Compatible avec Python 3.13 + vrai streaming progressif
"""
import boto3
import json
import os
import time
import base64
from typing import Optional
from uuid import uuid4
from decimal import Decimal

from fastapi import FastAPI, HTTPException, Header
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
import uvicorn
import io
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfReader

# Clients AWS
bedrock_client = boto3.client('bedrock-runtime', region_name='eu-west-3')
dynamodb = boto3.resource('dynamodb', region_name='eu-west-3')
kms_client = boto3.client('kms', region_name='eu-west-3')

# Configuration
MODEL_ID = 'eu.anthropic.claude-sonnet-4-5-20250929-v1:0'
DYNAMODB_TABLE = os.environ.get('DYNAMODB_TABLE', 'claude-serverless-prod-chat-history')
USAGE_TRACKING_TABLE = os.environ.get('USAGE_TRACKING_TABLE', 'claude-serverless-prod-usage-tracking')
KMS_KEY_ID = os.environ.get('KMS_KEY_ID')

# Prix des modèles Claude par 1000 tokens (en USD)
MODEL_PRICING = {
    'claude-haiku-4.5': {
        'input': 0.001,
        'output': 0.005,
        'cache_write': 0.00125,
        'cache_read': 0.0001
    },
    'claude-sonnet-4.5': {
        'input': 0.003,
        'output': 0.015,
        'cache_write': 0.00375,
        'cache_read': 0.0003
    },
    'claude-sonnet-4': {
        'input': 0.003,
        'output': 0.015,
        'cache_write': 0.00375,
        'cache_read': 0.0003
    }
}

# FastAPI app
# FastAPI app
app = FastAPI(title="Claude Chat API with Streaming")


# ============= Fonctions de chiffrement KMS =============

def encrypt_message(plaintext: str) -> str:
    """Chiffre un message avec KMS et retourne le résultat en base64"""
    if not KMS_KEY_ID:
        print("WARNING: KMS_KEY_ID not set, message not encrypted")
        return plaintext
    
    try:
        response = kms_client.encrypt(
            KeyId=KMS_KEY_ID,
            Plaintext=plaintext.encode('utf-8')
        )
        # Encoder le ciphertext en base64 pour le stocker dans DynamoDB
        return base64.b64encode(response['CiphertextBlob']).decode('utf-8')
    except Exception as e:
        print(f"ERROR encrypting message: {e}")
        raise HTTPException(status_code=500, detail=f"Encryption error: {str(e)}")


def decrypt_message(ciphertext_b64: str) -> str:
    """Déchiffre un message depuis base64 avec KMS"""
    if not KMS_KEY_ID:
        print("WARNING: KMS_KEY_ID not set, returning message as-is")
        return ciphertext_b64
    
    try:
        # Décoder depuis base64
        ciphertext_blob = base64.b64decode(ciphertext_b64)
        response = kms_client.decrypt(
            CiphertextBlob=ciphertext_blob
        )
        return response['Plaintext'].decode('utf-8')
    except Exception as e:
        print(f"ERROR decrypting message: {e}")
        # Si le déchiffrement échoue, retourner le texte brut (compatibilité anciennes données)
        return ciphertext_b64


def calculate_cost(input_tokens: int, output_tokens: int, model_name: str = 'claude-sonnet-4.5') -> float:
    """Calcule le coût en USD basé sur le nombre de tokens"""
    pricing = MODEL_PRICING.get(model_name, MODEL_PRICING['claude-sonnet-4.5'])
    
    # Coût = (input_tokens / 1000) * prix_input + (output_tokens / 1000) * prix_output
    cost = (input_tokens / 1000.0) * pricing['input'] + (output_tokens / 1000.0) * pricing['output']
    return round(cost, 6)  # Arrondir à 6 décimales


def update_usage_stats(user_id: str, input_tokens: int, output_tokens: int, cost: float):
    """Met à jour les statistiques d'utilisation dans DynamoDB (mensuel)"""
    try:
        from datetime import datetime, timedelta
        import time
        
        # Format du mois: YYYY-MM
        current_month = datetime.now().strftime('%Y-%m')
        
        # Convertir le coût en Decimal pour DynamoDB
        cost_decimal = Decimal(str(cost))
        
        # Calculer le TTL (1 an à partir de maintenant) en timestamp Unix
        ttl_date = datetime.now() + timedelta(days=365)
        ttl_timestamp = int(time.mktime(ttl_date.timetuple()))
        
        # Mise à jour mensuelle
        monthly_table = dynamodb.Table(USAGE_TRACKING_TABLE)
        monthly_table.update_item(
            Key={
                'user_id': user_id,
                'month': current_month
            },
            UpdateExpression='ADD input_tokens :input, output_tokens :output, cost_usd :cost SET #ttl = :ttl',
            ExpressionAttributeNames={
                '#ttl': 'ttl'
            },
            ExpressionAttributeValues={
                ':input': input_tokens,
                ':output': output_tokens,
                ':cost': cost_decimal,
                ':ttl': ttl_timestamp
            }
        )
        
        print(f"✅ Updated usage stats for {user_id} - Month: {current_month}, Input: {input_tokens}, Output: {output_tokens}, Cost: ${cost}, TTL: {ttl_date.strftime('%Y-%m-%d')}")
    except Exception as e:
        # Ne pas bloquer l'application si le suivi échoue
        print(f"⚠️ WARNING: Failed to update usage stats: {e}")


# ============= Models Pydantic =============

class FileData(BaseModel):
    fileName: str
    fileType: str
    fileContent: str  # base64


class ChatRequest(BaseModel):
    message: str
    conversationId: Optional[str] = None
    files: Optional[list[FileData]] = None
    # Rétro-compatibilité
    fileContents: Optional[list[str]] = None


def extract_user_id(authorization: Optional[str]) -> Optional[str]:
    """Extraire le user ID depuis le JWT"""
    if not authorization or not authorization.startswith('Bearer '):
        return None
    
    try:
        token = authorization.split(' ')[1]
        parts = token.split('.')
        if len(parts) != 3:
            return None
        
        # Décoder le payload JWT
        payload = parts[1]
        padding = 4 - (len(payload) % 4)
        if padding != 4:
            payload += '=' * padding
        
        decoded = base64.urlsafe_b64decode(payload)
        claims = json.loads(decoded)
        return claims.get('sub')
    except Exception as e:
        print(f"Error extracting user ID: {e}")
        return None


def extract_text_from_file(file_content_b64: str, file_type: str, file_name: str) -> str:
    """Extraire le texte d'un fichier selon son type"""
    try:
        # Décoder le base64
        file_bytes = base64.b64decode(file_content_b64)
        file_io = io.BytesIO(file_bytes)
        
        print(f"Extracting {file_name} ({file_type}), size: {len(file_bytes)} bytes")
        
        # PDF
        if file_type == 'application/pdf' or file_name.lower().endswith('.pdf'):
            reader = PdfReader(file_io)
            text = '\n'.join([page.extract_text() for page in reader.pages])
            print(f"PDF extracted: {len(text)} chars")
            return text
        
        # DOCX
        elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or file_name.lower().endswith('.docx'):
            doc = Document(file_io)
            text = '\n'.join([para.text for para in doc.paragraphs])
            print(f"DOCX extracted: {len(text)} chars, {len(doc.paragraphs)} paragraphs")
            return text
        
        # XLSX
        elif file_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' or file_name.lower().endswith('.xlsx'):
            wb = load_workbook(file_io, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n=== Feuille: {sheet_name} ===\n")
                for row in sheet.iter_rows(values_only=True):
                    row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip():
                        text_parts.append(row_text)
            return '\n'.join(text_parts)
        
        # PPTX
        elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or file_name.lower().endswith('.pptx'):
            prs = Presentation(file_io)
            text_parts = []
            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n=== Slide {i} ===\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return '\n'.join(text_parts)
        
        # TXT, CSV, JSON, etc.
        elif file_type.startswith('text/') or file_name.lower().endswith(('.txt', '.csv', '.json', '.md')):
            return file_bytes.decode('utf-8', errors='ignore')
        
        else:
            return f"[Fichier {file_name}: type non supporté pour extraction de texte]"
    
    except Exception as e:
        print(f"Error extracting text from file: {e}")
        return f"[Erreur lors de la lecture du fichier {file_name}: {str(e)}]"


def get_conversation_history(user_id: str, conversation_id: str) -> list:
    """Récupérer l'historique de conversation avec déchiffrement"""
    try:
        table = dynamodb.Table(DYNAMODB_TABLE)
        response = table.get_item(
            Key={
                'user_id': user_id,
                'conversation_id': conversation_id
            }
        )
        
        # Récupérer les messages chiffrés
        encrypted_messages = response.get('Item', {}).get('messages', [])
        
        # Déchiffrer le contenu de chaque message
        decrypted_messages = []
        for msg in encrypted_messages:
            decrypted_msg = msg.copy()
            # Déchiffrer uniquement le contenu
            decrypted_msg['content'] = decrypt_message(msg['content'])
            decrypted_messages.append(decrypted_msg)
        
        return decrypted_messages
    except Exception as e:
        print(f"Error getting conversation: {e}")
        return []


def save_conversation(user_id: str, conversation_id: str, messages: list):
    """Sauvegarder la conversation avec chiffrement des messages"""
    try:
        table = dynamodb.Table(DYNAMODB_TABLE)
        
        # Chiffrer le contenu de chaque message avant de sauvegarder
        encrypted_messages = []
        for msg in messages:
            encrypted_msg = msg.copy()
            # Chiffrer uniquement le contenu, pas les métadonnées
            encrypted_msg['content'] = encrypt_message(msg['content'])
            encrypted_messages.append(encrypted_msg)
        
        # Le TTL suffit pour limiter le nombre de conversations (90 jours)
        # Pas de limite artificielle sur le nombre de messages
        # (DynamoDB limite à 400KB par item, mais ça permet ~1000+ messages)
        
        ttl = int(time.time()) + (90 * 24 * 60 * 60)  # 90 jours
        
        table.put_item(
            Item={
                'user_id': user_id,
                'conversation_id': conversation_id,
                'messages': encrypted_messages,
                'timestamp': int(time.time() * 1000),
                'ttl': ttl
            }
        )
    except Exception as e:
        print(f"Error saving conversation: {e}")


async def stream_bedrock_response(
    messages: list,
    conversation_id: str,
    user_id: str,
    conversation_history: list,
    user_message: dict
):
    """Générateur asynchrone pour streamer depuis Bedrock"""
    
    # Envoyer métadonnées de début
    yield json.dumps({
        'type': 'start',
        'conversationId': conversation_id,
        'timestamp': int(time.time() * 1000)
    }) + '\n'
    
    # Formater les messages pour Bedrock
    formatted_messages = [
        {'role': msg['role'], 'content': msg['content']}
        for msg in messages
    ]
    
    request_body = {
        'anthropic_version': 'bedrock-2023-05-31',
        'max_tokens': 4000,
        'messages': formatted_messages,
        'system': 'Tu es un assistant IA utile et bienveillant. Tu peux analyser des documents et répondre aux questions à leur sujet. Réponds de manière claire et structurée.'
    }
    
    full_response = ''
    input_tokens = 0
    output_tokens = 0
    
    try:
        # Appel streaming à Bedrock
        response = bedrock_client.invoke_model_with_response_stream(
            modelId=MODEL_ID,
            contentType='application/json',
            body=json.dumps(request_body)
        )
        
        # Traiter le stream
        stream = response.get('body')
        if stream:
            for event in stream:
                chunk = event.get('chunk')
                if chunk:
                    chunk_data = json.loads(chunk.get('bytes').decode())
                    
                    # Debug: afficher le type de chunk reçu
                    print(f"🔍 Chunk type: {chunk_data.get('type')}, keys: {list(chunk_data.keys())}")
                    
                    if chunk_data['type'] == 'content_block_delta':
                        if 'delta' in chunk_data and 'text' in chunk_data['delta']:
                            text_chunk = chunk_data['delta']['text']
                            full_response += text_chunk
                            
                            # Envoyer le chunk au client
                            yield json.dumps({
                                'type': 'chunk',
                                'content': text_chunk
                            }) + '\n'
                    
                    # Capturer les statistiques de tokens
                    elif chunk_data['type'] == 'message_delta':
                        # Les tokens sont souvent dans message_delta
                        if 'usage' in chunk_data:
                            usage = chunk_data['usage']
                            # output_tokens est dans usage
                            if usage.get('output_tokens', 0) > 0:
                                output_tokens = usage.get('output_tokens', 0)
                                print(f"📊 Output tokens from message_delta: {output_tokens}")
                    
                    elif chunk_data['type'] == 'message_stop':
                        # Les tokens complets sont dans amazon-bedrock-invocationMetrics
                        if 'amazon-bedrock-invocationMetrics' in chunk_data:
                            metrics = chunk_data['amazon-bedrock-invocationMetrics']
                            input_tokens = metrics.get('inputTokenCount', 0)
                            output_tokens = metrics.get('outputTokenCount', 0)
                            print(f"� Complete token usage from invocationMetrics - Input: {input_tokens}, Output: {output_tokens}")
                    
                    # Vérifier si usage est directement dans le chunk (backup)
                    if 'usage' in chunk_data and input_tokens == 0:
                        usage = chunk_data['usage']
                        if usage.get('input_tokens', 0) > 0:
                            input_tokens = usage.get('input_tokens', 0)
                            print(f"📊 Input tokens from usage: {input_tokens}")
        
        # Calculer et enregistrer les coûts
        if input_tokens > 0 or output_tokens > 0:
            cost = calculate_cost(input_tokens, output_tokens)
            update_usage_stats(user_id, input_tokens, output_tokens, cost)
            print(f"💰 Cost for this request: ${cost}")
        
        # Envoyer métadonnées de fin
        yield json.dumps({
            'type': 'end',
            'timestamp': int(time.time() * 1000)
        }) + '\n'
        
        # Sauvegarder la conversation après streaming
        if full_response:
            assistant_message = {
                'role': 'assistant',
                'content': full_response,
                'timestamp': int(time.time() * 1000)
            }
            updated_messages = conversation_history + [user_message, assistant_message]
            save_conversation(user_id, conversation_id, updated_messages)
        
    except Exception as e:
        print(f"Error in Bedrock streaming: {e}")
        yield json.dumps({
            'type': 'error',
            'content': f'Error calling Claude: {str(e)}'
        }) + '\n'


@app.post("/chat")
async def chat_endpoint(
    request: ChatRequest,
    authorization: Optional[str] = Header(None)
):
    """Endpoint de chat avec streaming"""
    
    # Vérifier l'authentification
    user_id = extract_user_id(authorization)
    if not user_id:
        raise HTTPException(status_code=401, detail="Unauthorized")
    
    # Générer ou utiliser l'ID de conversation
    conversation_id = request.conversationId or str(uuid4())
    
    # Récupérer l'historique
    conversation_history = get_conversation_history(user_id, conversation_id)
    
    # Construire le contexte
    context_messages = conversation_history.copy()
    
    # Traiter les fichiers
    files_metadata = []
    files_text = []
    
    # Support du nouveau format avec métadonnées
    if request.files:
        for file in request.files:
            text = extract_text_from_file(file.fileContent, file.fileType, file.fileName)
            files_text.append(f"<fichier nom='{file.fileName}'>\n{text}\n</fichier>")
            files_metadata.append({
                'name': file.fileName,
                'type': file.fileType
            })
    # Rétro-compatibilité avec l'ancien format
    elif request.fileContents:
        for i, content in enumerate(request.fileContents):
            file_name = f"document_{i+1}.txt"
            text = extract_text_from_file(content, 'text/plain', file_name)
            files_text.append(f"<fichier nom='{file_name}'>\n{text}\n</fichier>")
            files_metadata.append({
                'name': file_name,
                'type': 'text/plain'
            })
    
    # Ajouter les fichiers au contexte si présents
    if files_text:
        context_messages.append({
            'role': 'user',
            'content': f"Voici les fichiers fournis en contexte:\n\n{chr(10).join(files_text)}",
            'timestamp': int(time.time() * 1000) - 1
        })
    
    # Ajouter le message utilisateur
    timestamp = int(time.time() * 1000)
    user_message = {
        'role': 'user',
        'content': request.message,
        'timestamp': timestamp,
        'files': files_metadata if files_metadata else None
    }
    context_messages.append(user_message)
    
    # Retourner le streaming response
    return StreamingResponse(
        stream_bedrock_response(
            context_messages,
            conversation_id,
            user_id,
            conversation_history,
            user_message
        ),
        media_type='application/x-ndjson',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no'  # Disable nginx buffering
        }
    )


@app.get("/conversations")
async def list_conversations_endpoint(
    authorization: Optional[str] = Header(None)
):
    """Lister toutes les conversations d'un utilisateur"""
    
    # Vérifier l'authentification
    user_id = extract_user_id(authorization)
    if not user_id:
        raise HTTPException(status_code=401, detail="Unauthorized")
    
    try:
        table = dynamodb.Table(DYNAMODB_TABLE)
        
        # Query DynamoDB pour récupérer toutes les conversations de l'utilisateur
        response = table.query(
            KeyConditionExpression='user_id = :user_id',
            ExpressionAttributeValues={
                ':user_id': user_id
            }
        )
        
        # Extraire les conversations et les trier par timestamp (plus récent en premier)
        conversations = []
        for item in response.get('Items', []):
            # Récupérer le premier et dernier message pour l'aperçu
            messages = item.get('messages', [])
            if messages:
                # Déchiffrer le premier message pour l'aperçu
                first_message_encrypted = messages[0].get('content', '')
                first_message = decrypt_message(first_message_encrypted)[:100]  # Premier message déchiffré et tronqué
                conversations.append({
                    'conversationId': item['conversation_id'],
                    'timestamp': item.get('timestamp', 0),
                    'messageCount': len(messages),
                    'preview': first_message
                })
        
        # Trier par timestamp décroissant (plus récent en premier)
        conversations.sort(key=lambda x: x['timestamp'], reverse=True)
        
        return {
            'conversations': conversations,
            'count': len(conversations)
        }
    
    except Exception as e:
        print(f"Error listing conversations: {e}")
        raise HTTPException(status_code=500, detail=f"Error listing conversations: {str(e)}")


@app.get("/conversations/{conversation_id}")
async def get_conversation_endpoint(
    conversation_id: str,
    authorization: Optional[str] = Header(None)
):
    """Récupérer l'historique d'une conversation"""
    
    # Vérifier l'authentification
    user_id = extract_user_id(authorization)
    if not user_id:
        raise HTTPException(status_code=401, detail="Unauthorized")
    
    # Récupérer l'historique
    messages = get_conversation_history(user_id, conversation_id)
    
    return {
        'conversationId': conversation_id,
        'messages': messages
    }


@app.delete("/conversations/{conversation_id}")
async def delete_conversation_endpoint(
    conversation_id: str,
    authorization: Optional[str] = Header(None)
):
    """Supprimer une conversation"""
    
    # Vérifier l'authentification
    user_id = extract_user_id(authorization)
    if not user_id:
        raise HTTPException(status_code=401, detail="Unauthorized")
    
    try:
        table = dynamodb.Table(DYNAMODB_TABLE)
        table.delete_item(
            Key={
                'user_id': user_id,
                'conversation_id': conversation_id
            }
        )
        return {'success': True, 'conversationId': conversation_id}
    except Exception as e:
        print(f"Error deleting conversation: {e}")
        raise HTTPException(status_code=500, detail=f"Error deleting conversation: {str(e)}")


@app.delete("/conversations")
async def delete_all_conversations_endpoint(
    authorization: Optional[str] = Header(None)
):
    """Supprimer TOUTES les conversations d'un utilisateur"""
    
    # Vérifier l'authentification
    user_id = extract_user_id(authorization)
    if not user_id:
        raise HTTPException(status_code=401, detail="Unauthorized")
    
    try:
        table = dynamodb.Table(DYNAMODB_TABLE)
        
        # Récupérer toutes les conversations de l'utilisateur
        response = table.query(
            KeyConditionExpression='user_id = :user_id',
            ExpressionAttributeValues={
                ':user_id': user_id
            }
        )
        
        conversations = response.get('Items', [])
        deleted_count = 0
        
        # Supprimer chaque conversation
        for conversation in conversations:
            table.delete_item(
                Key={
                    'user_id': user_id,
                    'conversation_id': conversation['conversation_id']
                }
            )
            deleted_count += 1
        
        print(f"Deleted {deleted_count} conversations for user {user_id}")
        return {
            'success': True, 
            'deletedCount': deleted_count,
            'message': f'{deleted_count} conversation(s) supprimée(s)'
        }
    except Exception as e:
        print(f"Error deleting all conversations: {e}")
        raise HTTPException(status_code=500, detail=f"Error deleting conversations: {str(e)}")


@app.get("/usage/current-month")
async def get_current_month_usage(
    authorization: Optional[str] = Header(None)
):
    """Récupérer les statistiques d'utilisation du mois en cours"""
    
    # Vérifier l'authentification
    user_id = extract_user_id(authorization)
    if not user_id:
        raise HTTPException(status_code=401, detail="Unauthorized")
    
    try:
        from datetime import datetime
        
        # Format du mois: YYYY-MM
        current_month = datetime.now().strftime('%Y-%m')
        
        table = dynamodb.Table(USAGE_TRACKING_TABLE)
        
        response = table.get_item(
            Key={
                'user_id': user_id,
                'month': current_month
            }
        )
        
        item = response.get('Item', {})
        
        return {
            'month': current_month,
            'inputTokens': int(item.get('input_tokens', 0)),
            'outputTokens': int(item.get('output_tokens', 0)),
            'costUsd': float(item.get('cost_usd', 0))
        }
    except Exception as e:
        print(f"Error getting usage stats: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting usage stats: {str(e)}")


@app.get("/usage/monthly-year")
async def get_monthly_usage_year(
    authorization: Optional[str] = Header(None)
):
    """Récupérer les statistiques mensuelles de l'année en cours"""
    
    # Vérifier l'authentification
    user_id = extract_user_id(authorization)
    if not user_id:
        raise HTTPException(status_code=401, detail="Unauthorized")
    
    try:
        from datetime import datetime
        
        now = datetime.now()
        current_year = now.year
        
        table = dynamodb.Table(USAGE_TRACKING_TABLE)
        
        # Query tous les mois de l'année pour cet utilisateur
        response = table.query(
            KeyConditionExpression='user_id = :user_id AND begins_with(#month, :year)',
            ExpressionAttributeNames={
                '#month': 'month'
            },
            ExpressionAttributeValues={
                ':user_id': user_id,
                ':year': str(current_year)
            }
        )
        
        items = response.get('Items', [])
        
        # Formatter les données
        monthly_data = []
        total_cost = Decimal('0')
        
        for item in items:
            cost = Decimal(str(item.get('cost_usd', 0)))
            total_cost += cost
            monthly_data.append({
                'month': item['month'],
                'inputTokens': int(item.get('input_tokens', 0)),
                'outputTokens': int(item.get('output_tokens', 0)),
                'costUsd': float(cost)
            })
        
        # Trier par mois
        monthly_data.sort(key=lambda x: x['month'])
        
        return {
            'year': current_year,
            'monthlyData': monthly_data,
            'totalCost': float(total_cost)
        }
    except Exception as e:
        print(f"Error getting monthly usage: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting monthly usage: {str(e)}")


@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {"status": "ok", "version": "1.0.0-lwa"}


if __name__ == "__main__":
    # Pour tests locaux
    uvicorn.run(app, host="0.0.0.0", port=int(os.environ.get("PORT", "8080")))
